from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.documents.pptx_parser import PptxParser


class UnavailableOfficeConverter:
    available = False


def test_pptx_parser_preserves_slide_text_image_order_and_context(tmp_path) -> None:
    image_buffer = BytesIO()
    Image.new("RGB", (80, 50), (239, 68, 68)).save(image_buffer, format="PNG")
    image_buffer.seek(0)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.5))
    title.text = "告警流程"
    before = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(5), Inches(0.5))
    before.text = "先检查服务状态。"
    slide.shapes.add_picture(image_buffer, Inches(0.5), Inches(2), width=Inches(2))
    after = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(5), Inches(0.5))
    after.text = "红色告警需要立即处理。"
    path = tmp_path / "alerts.pptx"
    presentation.save(path)

    elements = PptxParser(UnavailableOfficeConverter()).parse(path)
    image = next(element for element in elements if element.kind == "image")

    assert [element.kind for element in elements] == ["text", "text", "image", "text"]
    assert image.heading_path == ["第 1 页"]
    assert image.metadata["previous_text"] == "先检查服务状态。"
    assert image.metadata["next_text"] == "红色告警需要立即处理。"
    assert image.image_bytes


class RenderedOfficeConverter:
    available = True

    def __init__(self, pdf_path) -> None:
        self.pdf_path = pdf_path

    def convert(self, _source, _extension):
        class Conversion:
            def __enter__(inner_self):
                return self.pdf_path

            def __exit__(inner_self, *_args):
                return None

        return Conversion()


def test_pptx_full_slide_image_receives_all_slide_text_as_context(tmp_path) -> None:
    import pymupdf

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(0.5)).text = "第一步"
    slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(0.5)).text = "第二步"
    pptx_path = tmp_path / "steps.pptx"
    presentation.save(pptx_path)

    pdf_path = tmp_path / "steps.pdf"
    document = pymupdf.open()
    page = document.new_page(width=400, height=300)
    page.insert_text((40, 60), "Steps")
    document.save(pdf_path)
    document.close()

    elements = PptxParser(RenderedOfficeConverter(pdf_path)).parse(pptx_path)
    full_slide = next(element for element in elements if element.metadata.get("full_slide"))

    assert "第一步" in full_slide.metadata["previous_text"]
    assert "第二步" in full_slide.metadata["previous_text"]
