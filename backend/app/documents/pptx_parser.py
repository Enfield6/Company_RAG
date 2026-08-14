from pathlib import Path
from typing import Any

import pymupdf
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.documents.base import append_element, attach_image_neighborhood
from app.documents.office_converter import OfficeConverter
from app.documents.types import ParsedElement


class PptxParser:
    """Extract slide text, tables and pictures in top-to-bottom visual order."""

    def __init__(self, office_converter: OfficeConverter | None = None) -> None:
        self.office_converter = office_converter or OfficeConverter()

    def parse(self, path: Path) -> list[ParsedElement]:
        presentation = Presentation(path)
        rendered_slides = self._render_slides(path)
        elements: list[ParsedElement] = []
        for slide_index, slide in enumerate(presentation.slides):
            slide_number = slide_index + 1
            title_shape = slide.shapes.title
            title = self._shape_text(title_shape) if title_shape is not None else ""
            heading_path = [title] if title else [f"第 {slide_number} 页"]
            if title:
                append_element(
                    elements,
                    ParsedElement(
                        sequence_no=0,
                        kind="text",
                        text=title,
                        heading_path=heading_path,
                        metadata={
                            "source": "pptx",
                            "slide_number": slide_number,
                            "heading_level": 1,
                        },
                    ),
                )
            if slide_index < len(rendered_slides):
                append_element(
                    elements,
                    ParsedElement(
                        sequence_no=0,
                        kind="image",
                        heading_path=heading_path,
                        image_bytes=rendered_slides[slide_index],
                        image_extension=".png",
                        image_content_type="image/png",
                        relationship_id=f"slide-{slide_number}-render",
                        image_caption=f"第 {slide_number} 页完整幻灯片",
                        metadata={
                            "source": "pptx_slide_render",
                            "slide_number": slide_number,
                            "full_slide": True,
                        },
                    ),
                )
            shapes = sorted(
                self._flatten_shapes(slide.shapes),
                key=lambda shape: (int(getattr(shape, "top", 0)), int(getattr(shape, "left", 0))),
            )
            for shape_index, shape in enumerate(shapes):
                metadata = {
                    "source": "pptx",
                    "slide_number": slide_number,
                    "shape_index": shape_index,
                    "left": int(getattr(shape, "left", 0)),
                    "top": int(getattr(shape, "top", 0)),
                    "width": int(getattr(shape, "width", 0)),
                    "height": int(getattr(shape, "height", 0)),
                }
                if shape is title_shape:
                    continue
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    if slide_index < len(rendered_slides):
                        continue
                    image = shape.image
                    extension = f".{image.ext.lower()}"
                    append_element(
                        elements,
                        ParsedElement(
                            sequence_no=0,
                            kind="image",
                            heading_path=heading_path,
                            image_bytes=image.blob,
                            image_extension=extension,
                            image_content_type=image.content_type,
                            relationship_id=getattr(shape._element, "blip_rId", None),
                            image_caption=self._picture_caption(shape),
                            metadata=metadata,
                        ),
                    )
                    continue
                if getattr(shape, "has_table", False):
                    table_text = "\n".join(
                        " | ".join(cell.text.strip() for cell in row.cells)
                        for row in shape.table.rows
                    ).strip()
                    if table_text:
                        append_element(
                            elements,
                            ParsedElement(
                                sequence_no=0,
                                kind="table",
                                text=table_text,
                                heading_path=heading_path,
                                metadata=metadata,
                            ),
                        )
                    continue
                text = self._shape_text(shape)
                if text:
                    append_element(
                        elements,
                        ParsedElement(
                            sequence_no=0,
                            kind="text",
                            text=text,
                            heading_path=heading_path,
                            metadata=metadata,
                        ),
                    )
            notes_text_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            notes = notes_text_frame.text.strip() if notes_text_frame is not None else ""
            if notes:
                append_element(
                    elements,
                    ParsedElement(
                        sequence_no=0,
                        kind="text",
                        text=notes,
                        heading_path=heading_path,
                        metadata={"source": "pptx_notes", "slide_number": slide_number},
                    ),
                )
        attach_image_neighborhood(elements)
        slide_texts: dict[int, list[str]] = {}
        for element in elements:
            slide_number = element.metadata.get("slide_number")
            if isinstance(slide_number, int) and element.kind != "image" and element.text:
                slide_texts.setdefault(slide_number, []).append(element.text)
        for element in elements:
            if not element.metadata.get("full_slide"):
                continue
            slide_number = element.metadata.get("slide_number")
            context = "\n".join(slide_texts.get(slide_number, []))
            element.metadata["previous_text"] = context[:4000]
            element.metadata["next_text"] = ""
        return elements

    def _render_slides(self, path: Path) -> list[bytes]:
        if not self.office_converter.available:
            return []
        try:
            with self.office_converter.convert(path, ".pdf") as pdf_path:
                with pymupdf.open(pdf_path) as document:
                    return [
                        page.get_pixmap(matrix=pymupdf.Matrix(1.5, 1.5), alpha=False).tobytes(
                            "png"
                        )
                        for page in document
                    ]
        except (RuntimeError, pymupdf.FileDataError):
            return []

    @classmethod
    def _flatten_shapes(cls, shapes) -> list[Any]:
        flattened: list[Any] = []
        for shape in shapes:
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                flattened.extend(cls._flatten_shapes(shape.shapes))
            else:
                flattened.append(shape)
        return flattened

    @staticmethod
    def _shape_text(shape) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        return "\n".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text.strip()
        )

    @staticmethod
    def _picture_caption(shape) -> str:
        name = str(getattr(shape, "name", "")).strip()
        return "" if name.lower().startswith(("picture ", "图片 ")) else name
